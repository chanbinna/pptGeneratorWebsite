from flask import Flask, render_template, request, send_file, redirect, url_for
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
import os
import io

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        file = request.files.get("file")
        content = request.form.get("content")

        if file and file.filename:
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(filepath)
            pptx_filepath = create_ppt_from_text(filepath)
            os.remove(filepath)
            return send_file(pptx_filepath, as_attachment=True)

        elif content:
            pptx_data = create_ppt_from_textarea(content)
            return send_file(pptx_data, as_attachment=True, download_name="찬양ppt.pptx")

        return redirect(url_for("index"))

    return render_template("index.html")

def create_ppt_from_text(text_file):
    prs = Presentation()
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)

    with open(text_file, "r", encoding="utf-8") as file:
        paragraphs, titles = get_paragraphs(file.readlines())

    for paragraph, title in zip(paragraphs, titles):
        create_slide(prs, paragraph, title)

    output_path = os.path.join(app.config['UPLOAD_FOLDER'], "찬양ppt.pptx")
    prs.save(output_path)
    return output_path

def create_ppt_from_textarea(content):
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(14.29)

    lines = content.splitlines()
    paragraphs, titles = get_paragraphs(lines)

    for paragraph, title in zip(paragraphs, titles):
        create_slide(prs, paragraph, title)

    pptx_data = io.BytesIO()
    prs.save(pptx_data)
    pptx_data.seek(0)
    return pptx_data

def get_paragraphs(lines):
    paragraphs = []
    titles = []
    paragraph = []
    title = "Default Title"  # Default title if none found

    for line in lines:
        line = line.strip()
        if line.startswith("<") and line.endswith(">"):
            # Save the current paragraph and title if exists
            if paragraph:
                paragraphs.append(paragraph)
                titles.append(title)
                paragraph = []
            # Update the title for the new section
            title = line[1:-1]
        elif line:
            paragraph.append(line)
        else:
            if paragraph:
                paragraphs.append(paragraph)
                titles.append(title)
                paragraph = []

    # Append the last paragraph and title if remaining
    if paragraph:
        paragraphs.append(paragraph)
        titles.append(title)

    return paragraphs, titles

def create_slide(prs, content, title):
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img_path = "static/background2.png"

    # Background image
    background = slide.shapes.add_picture(img_path, Cm(0), Cm(0), width=prs.slide_width, height=prs.slide_height)
    slide.shapes._spTree.remove(background._element)
    slide.shapes._spTree.insert(2, background._element)

    # Top-left title
    top_left_box = slide.shapes.add_textbox(Cm(0.88), Cm(0.81), Cm(10), Cm(2))
    top_left_frame = top_left_box.text_frame
    top_left_frame.text = title
    top_left_paragraph = top_left_frame.paragraphs[0]
    top_left_paragraph.font.name = "Malgun Gothic"
    top_left_paragraph.font.bold = True
    top_left_paragraph.font.size = Pt(18)
    top_left_paragraph.font.color.rgb = RGBColor(0, 0, 0)
    top_left_paragraph.alignment = PP_ALIGN.LEFT

    # Split Korean and English lines
    korean_lines = []
    english_lines = []
    in_english_block = False
    english_block = []

    for line in content:
        line = line.strip()
        if line.startswith("{"):
            in_english_block = True
            english_block.append(line[1:].strip())
        elif line.endswith("}") and in_english_block:
            in_english_block = False
            english_block.append(line[:-1].strip())
            english_lines.extend(english_block)
            english_block = []
        elif in_english_block:
            english_block.append(line)
        else:
            korean_lines.append(line)

    # ✅ Main slide layout based on presence of English lines
    if english_lines:
        # Korean box
        kor_box = slide.shapes.add_textbox(Cm(4.79), Cm(0.91), Cm(24.58), Cm(11.39))
        kor_frame = kor_box.text_frame
        kor_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        for line in korean_lines:
            p = kor_frame.add_paragraph() if kor_frame.text else kor_frame.paragraphs[0]
            p.text = line
            p.font.name = "Malgun Gothic"
            p.font.size = Pt(54)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(0, 0, 0)

        # English box
        # 회색 투명 배경 이미지를 박스 위치에 삽입
        slide.shapes.add_picture("static/gray.png", Cm(1.18), Cm(11.37), Cm(31.72), Cm(7.25))

        # 그 위에 텍스트박스 삽입 (투명도 필요 없음)
        eng_box = slide.shapes.add_textbox(Cm(1.18), Cm(11.37), Cm(31.72), Cm(7.25))
        eng_frame = eng_box.text_frame
        eng_frame.margin_top = 0
        eng_frame.margin_bottom = 0
        eng_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        for line in english_lines:
            p = eng_frame.add_paragraph() if eng_frame.text else eng_frame.paragraphs[0]
            p.text = line
            p.font.name = "Malgun Gothic"
            p.font.size = Pt(40)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(0x25, 0x37, 0x61)

    else:
        # Only Korean
        title_box = slide.shapes.add_textbox(Cm(0), Cm((19.05 - 3) / 2), prs.slide_width, Cm(3))
        title_frame = title_box.text_frame
        title_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        for line in content:
            p = title_frame.add_paragraph() if title_frame.text else title_frame.paragraphs[0]
            p.text = line
            p.font.name = "Malgun Gothic"
            p.font.size = Pt(54)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RGBColor(0, 0, 0)

if __name__ == "__main__":
    os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
    app.run(debug=True)